from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


WORKSPACE_ROOT = Path(__file__).resolve().parents[2]
PITCH_DIR = WORKSPACE_ROOT / "pitch"
ASSETS_DIR = PITCH_DIR / "assets"
OUTPUT_PPTX = PITCH_DIR / "Southlake-Agentic-Synthetic-Data-Studio-Deck.pptx"

BACKGROUND = RGBColor(244, 247, 244)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(16, 42, 67)
MUTED = RGBColor(90, 113, 132)
TEAL = RGBColor(15, 124, 130)
TEAL_SOFT = RGBColor(223, 248, 236)
SAND = RGBColor(242, 231, 207)
CORAL = RGBColor(255, 139, 107)
ORANGE = RGBColor(255, 111, 0)
SLATE = RGBColor(30, 41, 59)
SLATE_SOFT = RGBColor(248, 250, 252)
GREEN_TEXT = RGBColor(9, 97, 73)
GOLD_TEXT = RGBColor(122, 82, 0)
RED_TEXT = RGBColor(154, 52, 18)


def add_soft_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND

    bubble_one = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.45),
        Inches(-0.55),
        Inches(2.9),
        Inches(2.9),
    )
    bubble_one.fill.solid()
    bubble_one.fill.fore_color.rgb = TEAL_SOFT
    bubble_one.line.fill.background()

    bubble_two = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.2),
        Inches(-0.4),
        Inches(2.1),
        Inches(2.1),
    )
    bubble_two.fill.solid()
    bubble_two.fill.fore_color.rgb = SAND
    bubble_two.line.fill.background()


def add_card(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.0)
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font_size=20,
    color=INK,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items,
    left,
    top,
    width,
    height,
    *,
    font_size=18,
    color=INK,
    level=0,
    bullet_color=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = level
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)
        paragraph.bullet = True
        if bullet_color is not None:
            paragraph.runs[0].font.color.rgb = bullet_color
    return box


def add_header(slide, eyebrow, title, subtitle, slide_number):
    return add_header_with_layout(slide, eyebrow, title, subtitle, slide_number)


def add_header_with_layout(
    slide,
    eyebrow,
    title,
    subtitle,
    slide_number,
    *,
    title_font_size=26,
    title_height=1.18,
    subtitle_top=1.72,
    subtitle_height=0.34,
):
    add_text(
        slide,
        eyebrow.upper(),
        Inches(0.72),
        Inches(0.42),
        Inches(5.0),
        Inches(0.35),
        font_size=11,
        color=TEAL,
        bold=True,
        font_name="Aptos",
    )
    add_text(
        slide,
        title,
        Inches(0.72),
        Inches(0.72),
        Inches(8.8),
        Inches(title_height),
        font_size=title_font_size,
        color=INK,
        bold=True,
        font_name="Aptos Display",
    )
    add_text(
        slide,
        subtitle,
        Inches(0.72),
        Inches(subtitle_top),
        Inches(8.9),
        Inches(subtitle_height),
        font_size=13,
        color=MUTED,
    )
    slide_badge = add_card(
        slide,
        Inches(11.8),
        Inches(0.42),
        Inches(0.8),
        Inches(0.38),
        fill_rgb=WHITE,
        line_rgb=RGBColor(220, 227, 234),
    )
    badge_frame = slide_badge.text_frame
    badge_frame.clear()
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = badge_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"{slide_number}/3"
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = MUTED


def add_footer(slide, text):
    add_text(
        slide,
        text,
        Inches(0.72),
        Inches(7.02),
        Inches(9.5),
        Inches(0.25),
        font_size=10,
        color=MUTED,
    )


def build_slide_one(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_header_with_layout(
        slide,
        "Methodology",
        "From Public ED Data To Southlake Future-State Planning Scenarios",
        "Agentic synthetic data workflow for distributed-health-network planning questions.",
        1,
        title_font_size=24,
        title_height=1.22,
        subtitle_top=1.8,
        subtitle_height=0.3,
    )

    message_card = add_card(
        slide,
        Inches(0.62),
        Inches(2.18),
        Inches(4.55),
        Inches(1.68),
        fill_rgb=WHITE,
        line_rgb=RGBColor(228, 233, 237),
    )
    frame = message_card.text_frame
    frame.clear()
    frame.word_wrap = True
    title_paragraph = frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = "Core Message"
    title_run.font.name = "Aptos"
    title_run.font.size = Pt(11)
    title_run.font.bold = True
    title_run.font.color.rgb = TEAL

    body_paragraph = frame.add_paragraph()
    body_paragraph.text = (
        "We turn safe public ED data into synthetic planning artifacts for "
        "Southlake's distributed-health-network questions, then make the "
        "limitations visible before anyone over-trusts the output."
    )
    body_paragraph.font.name = "Aptos"
    body_paragraph.font.size = Pt(17)
    body_paragraph.font.bold = True
    body_paragraph.font.color.rgb = INK
    body_paragraph.space_after = Pt(12)

    method_cards = [
        (
            "Input",
            "Curated public ED dataset used as a safe planning stand-in.",
            TEAL_SOFT,
            GREEN_TEXT,
        ),
        (
            "Agent Flow",
            "Goal framing -> profile -> plan -> generate -> evaluate -> caution report.",
            WHITE,
            INK,
        ),
        (
            "Southlake Lens",
            "Distributed campus routing, community diversion, growth, and surge scenarios.",
            SAND,
            GOLD_TEXT,
        ),
        (
            "Outputs",
            "Synthetic dataset, metrics dashboard, caution summary, export ZIP, pitch artifacts.",
            WHITE,
            INK,
        ),
    ]
    positions = [
        (Inches(0.62), Inches(4.06)),
        (Inches(2.95), Inches(4.06)),
        (Inches(0.62), Inches(5.5)),
        (Inches(2.95), Inches(5.5)),
    ]
    for (title, body, fill_rgb, heading_rgb), (left, top) in zip(method_cards, positions):
        card = add_card(
            slide,
            left,
            top,
            Inches(2.22),
            Inches(1.28),
            fill_rgb=fill_rgb,
            line_rgb=RGBColor(228, 233, 237) if fill_rgb == WHITE else None,
        )
        card.text_frame.clear()
        card.text_frame.word_wrap = True
        heading = card.text_frame.paragraphs[0]
        heading_run = heading.add_run()
        heading_run.text = title
        heading_run.font.name = "Aptos"
        heading_run.font.size = Pt(12)
        heading_run.font.bold = True
        heading_run.font.color.rgb = heading_rgb
        paragraph = card.text_frame.add_paragraph()
        paragraph.text = body
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = INK
        paragraph.space_before = Pt(3)
        paragraph.space_after = Pt(0)

    graphic_card = add_card(
        slide,
        Inches(5.48),
        Inches(2.02),
        Inches(7.25),
        Inches(4.7),
        fill_rgb=WHITE,
        line_rgb=RGBColor(228, 233, 237),
    )
    graphic_card.shadow.inherit = False
    slide.shapes.add_picture(
        str(ASSETS_DIR / "architecture-workflow.png"),
        Inches(5.72),
        Inches(2.26),
        width=Inches(6.8),
        height=Inches(3.82),
    )

    ribbon = add_card(
        slide,
        Inches(5.48),
        Inches(6.22),
        Inches(7.25),
        Inches(0.48),
        fill_rgb=SLATE,
    )
    ribbon_frame = ribbon.text_frame
    ribbon_frame.clear()
    ribbon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = ribbon_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "Planning questions: campus routing | observation demand | transfer pathways | closer-to-home care"
    run.font.name = "Aptos"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_footer(
        slide,
        "Message for judges: the product is not only about generating rows; it is about making synthetic planning data usable and governed.",
    )


def build_slide_two(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_header_with_layout(
        slide,
        "Features",
        "What The Studio Actually Does",
        "Saved demo reliability, visible evaluation, and exportable artifacts for innovation teams.",
        2,
        title_font_size=25,
        title_height=0.82,
        subtitle_top=1.55,
        subtitle_height=0.26,
    )

    shot_card = add_card(
        slide,
        Inches(0.62),
        Inches(1.82),
        Inches(7.2),
        Inches(5.3),
        fill_rgb=WHITE,
        line_rgb=RGBColor(228, 233, 237),
    )
    slide.shapes.add_picture(
        str(ASSETS_DIR / "run-focus.png"),
        Inches(0.84),
        Inches(2.04),
        width=Inches(6.76),
        height=Inches(4.51),
    )

    labels = [
        ("Saved demo run", Inches(0.98), Inches(1.92), TEAL_SOFT, GREEN_TEXT),
        ("Scenario controls", Inches(2.68), Inches(1.92), SAND, GOLD_TEXT),
        ("Metrics", Inches(4.36), Inches(1.92), RGBColor(231, 242, 255), TEAL),
        ("Pitch + cautions", Inches(5.53), Inches(6.22), RGBColor(255, 239, 234), RED_TEXT),
    ]
    for text, left, top, fill_rgb, font_rgb in labels:
        chip = add_card(
            slide,
            left,
            top,
            Inches(1.44),
            Inches(0.34),
            fill_rgb=fill_rgb,
            line_rgb=None,
        )
        chip.text_frame.clear()
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = chip.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = font_rgb

    add_text(
        slide,
        "Feature Highlights",
        Inches(8.1),
        Inches(1.92),
        Inches(3.9),
        Inches(0.35),
        font_size=12,
        color=TEAL,
        bold=True,
    )
    add_bullets(
        slide,
        [
            "Southlake-specific scenarios aligned to distributed-network planning needs",
            "One-click saved demo runs plus optional fresh synthesis or CSV upload",
            "Visible fidelity, privacy, and utility metrics in one screen",
            "Auto-generated methodology, features, and cautions for pitch preparation",
        ],
        Inches(8.1),
        Inches(2.2),
        Inches(4.0),
        Inches(2.45),
        font_size=14,
    )

    saved_run_card = add_card(
        slide,
        Inches(8.1),
        Inches(4.92),
        Inches(4.0),
        Inches(1.52),
        fill_rgb=SLATE,
    )
    saved_run_card.text_frame.clear()
    saved_run_card.text_frame.word_wrap = True
    saved_run_card.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    heading = saved_run_card.text_frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = "Primary saved run"
    heading_run.font.name = "Aptos"
    heading_run.font.size = Pt(11)
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor(174, 239, 215)

    run_line = saved_run_card.text_frame.add_paragraph()
    run_line.text = "Distributed Campus Routing | run 72c9c81912"
    run_line.font.name = "Aptos"
    run_line.font.size = Pt(16)
    run_line.font.bold = True
    run_line.font.color.rgb = WHITE
    run_line.space_after = Pt(6)

    metric_line = saved_run_card.text_frame.add_paragraph()
    metric_line.text = "Fidelity 81.54   Privacy 100.0   Utility 88.92"
    metric_line.font.name = "Aptos"
    metric_line.font.size = Pt(13)
    metric_line.font.color.rgb = WHITE

    add_footer(
        slide,
        "Pitch note: judges can see both the data product and the explanation product, with a reliable saved-run path if live generation is slow.",
    )


def build_slide_three(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_soft_background(slide)
    add_header_with_layout(
        slide,
        "Cautions",
        slide_number=3,
        title="What The Studio Is Safe For, And What It Is Not",
        subtitle="Synthetic data is useful for planning and experimentation, but it is not local operational truth.",
        title_font_size=22,
        title_height=1.05,
        subtitle_top=1.76,
        subtitle_height=0.26,
    )

    left_col = add_card(
        slide,
        Inches(0.62),
        Inches(2.18),
        Inches(5.86),
        Inches(4.54),
        fill_rgb=TEAL_SOFT,
        line_rgb=None,
    )
    left_col.text_frame.clear()
    left_col.text_frame.word_wrap = True
    heading = left_col.text_frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = "Good For"
    heading_run.font.name = "Aptos Display"
    heading_run.font.size = Pt(24)
    heading_run.font.bold = True
    heading_run.font.color.rgb = GREEN_TEXT
    for item in [
        "Planning, innovation, and prototyping conversations",
        "Scenario workshops about distributed-campus routing and community diversion",
        "Testing dashboards, workflows, and early-stage analytics without real patient records",
        "Pitch-day storytelling backed by metrics, cautions, and exportable artifacts",
    ]:
        paragraph = left_col.text_frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(10)

    right_col = add_card(
        slide,
        Inches(6.72),
        Inches(2.18),
        Inches(5.98),
        Inches(4.54),
        fill_rgb=RGBColor(255, 239, 234),
        line_rgb=None,
    )
    right_col.text_frame.clear()
    right_col.text_frame.word_wrap = True
    heading = right_col.text_frame.paragraphs[0]
    heading_run = heading.add_run()
    heading_run.text = "Not Ready For"
    heading_run.font.name = "Aptos Display"
    heading_run.font.size = Pt(24)
    heading_run.font.bold = True
    heading_run.font.color.rgb = RED_TEXT
    for item in [
        "Clinical care, staffing commitments, reimbursement, or patient-level action",
        "Claims about Southlake's real patient mix or operational truth from a U.S. public-use source",
        "A full hospital digital twin or a validated multi-table operational model",
        "Production use without governance, privacy review, and validation on governed local data",
    ]:
        paragraph = right_col.text_frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(10)

    bottom_bar = add_card(
        slide,
        Inches(0.92),
        Inches(6.16),
        Inches(11.6),
        Inches(0.58),
        fill_rgb=SLATE,
        line_rgb=None,
    )
    bottom_bar.text_frame.clear()
    bottom_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = bottom_bar.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "Best next step: validate the workflow with governed Southlake data, operational owners, and formal privacy review."
    run.font.name = "Aptos"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_footer(
        slide,
        "Credibility note: this caution slide is a strength. It shows the team understands healthcare risk and is not overselling the prototype.",
    )


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_one(prs)
    build_slide_two(prs)
    build_slide_three(prs)

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    path = build_deck()
    print(path)
